import os
import base64
import tempfile
import re
from flask import Flask, render_template, request, jsonify
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024
app.config['UPLOAD_FOLDER'] = tempfile.gettempdir()

ALLOWED_EXTENSIONS = {'pdf', 'pptx', 'docx'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def clean_text(text):
    """Clean extracted text"""
    if not text:
        return ""
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text)
    # Remove special characters but keep basic punctuation
    text = re.sub(r'[^\x00-\x7F]+', ' ', text)
    return text.strip()

def extract_pdf(filepath):
    """Extract text and images from PDF with page structure"""
    try:
        import fitz
        doc = fitz.open(filepath)
        slides_data = []
        all_images = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            clean = clean_text(text)
            
            if clean:  # Only add non-empty pages
                slides_data.append({
                    'number': page_num + 1,
                    'type': 'page',
                    'title': f"Page {page_num + 1}",
                    'content': clean,
                    'has_content': len(clean) > 50
                })
            
            # Extract images
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    ext = base_image["ext"]
                    b64 = base64.b64encode(image_bytes).decode('utf-8')
                    all_images.append({
                        "id": f"pdf_p{page_num}_img{img_index}",
                        "slide": page_num + 1,
                        "src": f"data:image/{ext};base64,{b64}",
                        "type": "pdf"
                    })
                except:
                    continue
        doc.close()
        return slides_data, all_images
        
    except ImportError:
        # Fallback
        text_content = []
        try:
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    clean = clean_text(text)
                    if clean:
                        text_content.append({
                            'number': i + 1,
                            'type': 'page',
                            'title': f"Page {i + 1}",
                            'content': clean,
                            'has_content': len(clean) > 50
                        })
        except:
            # Last resort - read as text
            with open(filepath, 'rb') as f:
                text = f.read().decode('utf-8', errors='ignore')
                text_content.append({
                    'number': 1,
                    'type': 'page',
                    'title': 'Document',
                    'content': clean_text(text),
                    'has_content': True
                })
        
        return text_content, []

def extract_pptx(filepath):
    """Extract slides with titles and content separately"""
    from pptx import Presentation
    
    prs = Presentation(filepath)
    slides_data = []
    images = []
    
    for idx, slide in enumerate(prs.slides, 1):
        slide_text = []
        slide_title = f"Slide {idx}"
        has_title = False
        
        # Extract title first
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Usually title is first shape or has larger font
                if not has_title and len(shape.text.strip()) < 100 and idx == 1:
                    slide_title = shape.text.strip()
                    has_title = True
                else:
                    slide_text.append(shape.text.strip())
            
            # Extract images
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext if hasattr(image, 'ext') else 'png'
                    b64 = base64.b64encode(image_bytes).decode('utf-8')
                    images.append({
                        "id": f"pptx_s{idx}_img{len(images)}",
                        "slide": idx,
                        "src": f"data:image/{ext};base64,{b64}",
                        "type": "pptx"
                    })
                except:
                    continue
        
        content = clean_text("\n".join(slide_text))
        
        # Try to infer title from first line if short
        lines = content.split('.')
        if len(lines) > 0 and len(lines[0]) < 80 and not has_title:
            slide_title = lines[0].strip()
            content = '.'.join(lines[1:]).strip()
        
        slides_data.append({
            'number': idx,
            'type': 'slide',
            'title': slide_title,
            'content': content,
            'has_content': len(content) > 30
        })
    
    return slides_data, images

def extract_docx(filepath):
    """Extract docx with section awareness"""
    from docx import Document
    
    doc = Document(filepath)
    sections = []
    images = []
    current_text = []
    section_count = 1
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
            
        # Check if it's a heading (simplistic check)
        if len(text) < 100 and (text.isupper() or text.endswith(':') or para.style.name.startswith('Heading')):
            if current_text:
                sections.append({
                    'number': section_count,
                    'type': 'section',
                    'title': f"Section {section_count}",
                    'content': clean_text("\n".join(current_text)),
                    'has_content': True
                })
                section_count += 1
                current_text = []
            sections.append({
                'number': section_count,
                'type': 'heading',
                'title': text,
                'content': '',
                'has_content': False
            })
        else:
            current_text.append(text)
    
    # Add remaining text
    if current_text:
        sections.append({
            'number': section_count,
            'type': 'section',
            'title': f"Section {section_count}",
            'content': clean_text("\n".join(current_text)),
            'has_content': True
        })
    
    # Extract images
    try:
        rels = doc.part.rels
        for rel in rels.values():
            if "image" in rel.target_ref:
                try:
                    image = rel.target_part.blob
                    ext = 'png'
                    if image[:4] == b'\x89PNG':
                        ext = 'png'
                    elif image[:2] == b'\xff\xd8':
                        ext = 'jpeg'
                    
                    b64 = base64.b64encode(image).decode('utf-8')
                    images.append({
                        "id": f"docx_img{len(images)}",
                        "slide": 1,
                        "src": f"data:image/{ext};base64,{b64}",
                        "type": "docx"
                    })
                except:
                    continue
    except:
        pass
    
    # If no sections found, treat as single document
    if not sections:
        full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        sections = [{
            'number': 1,
            'type': 'document',
            'title': 'Document',
            'content': clean_text(full_text),
            'has_content': True
        }]
    
    return sections, images

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/process', methods=['POST'])
def process_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file type. Use PDF, PPTX, or DOCX'}), 400
    
    try:
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        ext = filename.rsplit('.', 1)[1].lower()
        file_type = "document"
        
        if ext == 'pdf':
            slides_data, images = extract_pdf(filepath)
            file_type = "pdf"
        elif ext == 'pptx':
            slides_data, images = extract_pptx(filepath)
            file_type = "slides"
        elif ext == 'docx':
            slides_data, images = extract_docx(filepath)
            file_type = "document"
        else:
            return jsonify({'error': 'Unsupported extension'}), 400
        
        # Flatten text for backward compatibility
        full_text = "\n\n".join([f"## {s['title']}\n{s['content']}" for s in slides_data if s['has_content']])
        
        # Cleanup
        try:
            os.remove(filepath)
        except:
            pass
        
        return jsonify({
            'success': True,
            'slides': slides_data,  # New structured format
            'text': full_text,      # Legacy flat format
            'images': images,
            'file_type': file_type,
            'slide_count': len(slides_data)
        })
        
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    print("Starting UniNotes AI Server...")
    print("Open http://localhost:5000")
    app.run(debug=True, host='0.0.0.0', port=5000)